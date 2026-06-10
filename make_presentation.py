 rows"""Generate Depot-MVP PowerPoint presentation."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand colours ────────────────────────────────────────────────────────────
PURPLE      = RGBColor(0x76, 0x4B, 0xA2)
BLUE        = RGBColor(0x66, 0x7E, 0xEA)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK        = RGBColor(0x22, 0x22, 0x33)
LIGHT_BG    = RGBColor(0xF7, 0xF9, 0xFC)
GREEN       = RGBColor(0x4C, 0xAF, 0x50)
YELLOW      = RGBColor(0xFF, 0xC1, 0x07)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def prs() -> Presentation:
    p = Presentation()
    p.slide_width  = SLIDE_W
    p.slide_height = SLIDE_H
    return p


def blank(p: Presentation):
    return p.slides.add_slide(p.slide_layouts[6])  # completely blank


def bg(slide, color: RGBColor):
    from pptx.util import Emu
    from pptx.oxml.ns import qn
    from lxml import etree
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, left, top, width, height,
        fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def txt(slide, text, left, top, width, height,
        size=Pt(18), bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def bullet_box(slide, items, left, top, width, height,
               size=Pt(16), color=DARK, bullet="●"):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = size
        run.font.color.rgb = color
    return txb


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 – Title
# ─────────────────────────────────────────────────────────────────────────────
def slide_title(p):
    s = blank(p)
    bg(s, DARK)

    # Gradient accent bar (left)
    box(s, Inches(0), Inches(0), Inches(0.18), SLIDE_H, fill_color=BLUE)

    # Large title
    txt(s, "Depot-MVP",
        Inches(0.5), Inches(1.6), Inches(8), Inches(1.4),
        size=Pt(60), bold=True, color=WHITE)

    # Subtitle
    txt(s, "Automatische Portfolio-Analyse aus CSV & PDF",
        Inches(0.5), Inches(3.0), Inches(9), Inches(0.7),
        size=Pt(26), color=RGBColor(0xBB, 0xBB, 0xDD))

    # Tag line
    txt(s, "Interactive Brokers · Echtzeit-Auswertung · KI-Kommentar",
        Inches(0.5), Inches(3.9), Inches(9), Inches(0.5),
        size=Pt(18), color=BLUE)

    txt(s, "Tim Solovyev  ·  Dr. Mayer Consulting GmbH",
        Inches(0.5), Inches(6.5), Inches(9), Inches(0.5),
        size=Pt(15), color=RGBColor(0x88, 0x88, 0xAA))

    # Right visual block
    box(s, Inches(10.2), Inches(1.5), Inches(2.7), Inches(4.5),
        fill_color=RGBColor(0x33, 0x33, 0x55))
    txt(s, "📊", Inches(10.5), Inches(2.8), Inches(2), Inches(1.5),
        size=Pt(72), color=WHITE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 – Problem / Opportunity
# ─────────────────────────────────────────────────────────────────────────────
def slide_problem(p):
    s = blank(p)
    bg(s, LIGHT_BG)
    box(s, Inches(0), Inches(0), SLIDE_W, Inches(1.3), fill_color=PURPLE)
    txt(s, "Das Problem", Inches(0.5), Inches(0.25), Inches(10), Inches(0.8),
        size=Pt(32), bold=True, color=WHITE)

    problems = [
        "IB-Kontoauszüge als CSV/PDF sind unlesbar für Nicht-Techniker",
        "Manuelle Auswertung in Excel dauert Stunden",
        "Klumpenrisiken & Konzentration werden leicht übersehen",
        "Keine schnelle Übersicht über Asset-Mix und Gewichtung",
    ]
    bullet_box(s, problems,
               Inches(0.8), Inches(1.7), Inches(11.5), Inches(4.5),
               size=Pt(22), color=DARK)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 – Solution
# ─────────────────────────────────────────────────────────────────────────────
def slide_solution(p):
    s = blank(p)
    bg(s, DARK)
    box(s, Inches(0), Inches(0), SLIDE_W, Inches(1.3), fill_color=BLUE)
    txt(s, "Die Lösung", Inches(0.5), Inches(0.25), Inches(10), Inches(0.8),
        size=Pt(32), bold=True, color=WHITE)

    txt(s, "Datei hochladen  →  sofortige Analyse",
        Inches(0.5), Inches(1.5), Inches(12), Inches(0.7),
        size=Pt(28), bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    steps = [
        ("1", "CSV oder PDF hochladen", "Drag & Drop im Browser"),
        ("2", "Parser erkennt Format", "IB-Statement automatisch erkannt"),
        ("3", "Analyse in Sekunden",   "Wert, Gewichtung, Top-5, Risiken"),
        ("4", "KI-Kommentar",          "Klarer Text, keine Fachsprache"),
    ]
    col_w = Inches(2.9)
    for i, (num, title, sub) in enumerate(steps):
        x = Inches(0.4) + i * (col_w + Inches(0.3))
        box(s, x, Inches(2.5), col_w, Inches(3.5),
            fill_color=RGBColor(0x33, 0x33, 0x55))
        txt(s, num,  x + Inches(0.15), Inches(2.65), col_w, Inches(0.8),
            size=Pt(36), bold=True, color=BLUE, align=PP_ALIGN.CENTER)
        txt(s, title, x + Inches(0.1), Inches(3.5), col_w - Inches(0.2), Inches(0.7),
            size=Pt(16), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, sub,   x + Inches(0.1), Inches(4.25), col_w - Inches(0.2), Inches(1.2),
            size=Pt(13), color=RGBColor(0xAA, 0xAA, 0xCC), align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 – Features: Strukturierte Ausgabe
# ─────────────────────────────────────────────────────────────────────────────
def slide_output(p):
    s = blank(p)
    bg(s, LIGHT_BG)
    box(s, Inches(0), Inches(0), SLIDE_W, Inches(1.3), fill_color=PURPLE)
    txt(s, "Strukturierte Ausgabe", Inches(0.5), Inches(0.25), Inches(10), Inches(0.8),
        size=Pt(32), bold=True, color=WHITE)

    cols = [
        ("Wertpapier & Ticker", ["AAPL", "MSFT", "NVDA …"]),
        ("Stückzahl",           ["100", "50", "30 …"]),
        ("Aktueller Wert (EUR)",["17.454 €", "17.231 €", "24.183 €"]),
        ("Gewichtung %",        ["19,0 %", "18,8 %", "26,4 %"]),
        ("Asset-Typ",           ["Aktie", "ETF", "Option", "Anleihe", "Rohstoff", "Cash"]),
    ]

    col_w = Inches(2.3)
    for i, (header, values) in enumerate(cols):
        x = Inches(0.3) + i * (col_w + Inches(0.18))
        box(s, x, Inches(1.5), col_w, Inches(0.6), fill_color=BLUE)
        txt(s, header, x + Inches(0.1), Inches(1.55), col_w - Inches(0.1), Inches(0.5),
            size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        for j, v in enumerate(values):
            c = LIGHT_BG if j % 2 == 0 else WHITE
            box(s, x, Inches(2.1) + j * Inches(0.55), col_w, Inches(0.52),
                fill_color=RGBColor(0xEE, 0xF2, 0xFF) if j % 2 == 0 else WHITE)
            txt(s, v, x + Inches(0.1), Inches(2.13) + j * Inches(0.55),
                col_w - Inches(0.1), Inches(0.45),
                size=Pt(13), color=DARK, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 – Features: Analyse
# ─────────────────────────────────────────────────────────────────────────────
def slide_analysis(p):
    s = blank(p)
    bg(s, DARK)
    box(s, Inches(0), Inches(0), SLIDE_W, Inches(1.3), fill_color=BLUE)
    txt(s, "Portfolio-Analyse", Inches(0.5), Inches(0.25), Inches(10), Inches(0.8),
        size=Pt(32), bold=True, color=WHITE)

    cards = [
        (BLUE,   "💰 Gesamtwert",        "Summe aller Positionen\nin EUR"),
        (PURPLE, "🏆 Top-5 Positionen",  "Größte Positionen nach\nMarktwert & Gewichtung"),
        (GREEN,  "⚠️ Klumpenrisiken",    "Warnung ab 15 % Einzel-\nposition oder 40 % Top-5"),
        (YELLOW, "🌍 Tech/USA-Anteil",   "Erkennt Konzentration in\nUS-Technologie-Werte"),
    ]

    cw = Inches(2.9)
    for i, (color, title, desc) in enumerate(cards):
        x = Inches(0.4) + i * (cw + Inches(0.28))
        box(s, x, Inches(1.6), cw, Inches(4.8),
            fill_color=RGBColor(0x2A, 0x2A, 0x44))
        box(s, x, Inches(1.6), cw, Inches(0.12), fill_color=color)
        txt(s, title, x + Inches(0.15), Inches(1.8), cw - Inches(0.2), Inches(0.8),
            size=Pt(17), bold=True, color=WHITE)
        txt(s, desc,  x + Inches(0.15), Inches(2.65), cw - Inches(0.2), Inches(1.8),
            size=Pt(14), color=RGBColor(0xCC, 0xCC, 0xDD))


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 – Tech Stack
# ─────────────────────────────────────────────────────────────────────────────
def slide_tech(p):
    s = blank(p)
    bg(s, LIGHT_BG)
    box(s, Inches(0), Inches(0), SLIDE_W, Inches(1.3), fill_color=PURPLE)
    txt(s, "Tech-Stack", Inches(0.5), Inches(0.25), Inches(10), Inches(0.8),
        size=Pt(32), bold=True, color=WHITE)

    layers = [
        ("Frontend",  "HTML · CSS · Vanilla JS",       "Drag & Drop Upload, responsive Tabellen"),
        ("Backend",   "Python · Flask · Gunicorn",     "REST-API, Datei-Upload, JSON-Response"),
        ("Parser",    "pdfplumber · csv · pandas",     "IB-CSV Erkennung, FX-Umrechnung, PDF-Fallback"),
        ("Analyse",   "pandas · NumPy",                "Gewichtung, Top-N, Asset-Typ, Konzentration"),
        ("Deploy",    "AWS Elastic Beanstalk",         "Öffentlich erreichbar, automatisch skalierend"),
    ]

    for i, (layer, tech, desc) in enumerate(layers):
        y = Inches(1.6) + i * Inches(1.0)
        box(s, Inches(0.4), y, Inches(2.0), Inches(0.75), fill_color=BLUE)
        txt(s, layer, Inches(0.5), y + Inches(0.1), Inches(1.9), Inches(0.6),
            size=Pt(15), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, tech,  Inches(2.6), y + Inches(0.05), Inches(3.5), Inches(0.6),
            size=Pt(15), bold=True, color=DARK)
        txt(s, desc,  Inches(6.2), y + Inches(0.05), Inches(6.8), Inches(0.6),
            size=Pt(14), color=RGBColor(0x55, 0x55, 0x77))


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 – Demo / How to use
# ─────────────────────────────────────────────────────────────────────────────
def slide_demo(p):
    s = blank(p)
    bg(s, DARK)
    box(s, Inches(0), Inches(0), SLIDE_W, Inches(1.3), fill_color=GREEN)
    txt(s, "So funktioniert es", Inches(0.5), Inches(0.25), Inches(10), Inches(0.8),
        size=Pt(32), bold=True, color=WHITE)

    steps = [
        ("🖥",  "App öffnen",         "localhost:8080 im Browser aufrufen"),
        ("📄",  "Datei hochladen",    "IB-CSV oder PDF per Drag & Drop"),
        ("⚙️",  "Automatisch analysiert", "Parser erkennt Format, berechnet Kennzahlen"),
        ("📊",  "Ergebnisse lesen",   "Tabellen, Warnungen & KI-Kommentar"),
    ]

    for i, (icon, title, desc) in enumerate(steps):
        y = Inches(1.7) + i * Inches(1.3)
        box(s, Inches(0.5), y, Inches(0.9), Inches(0.9),
            fill_color=RGBColor(0x33, 0x33, 0x55))
        txt(s, icon, Inches(0.55), y + Inches(0.05), Inches(0.8), Inches(0.8),
            size=Pt(28), align=PP_ALIGN.CENTER)
        txt(s, title, Inches(1.6), y + Inches(0.0), Inches(4), Inches(0.5),
            size=Pt(18), bold=True, color=WHITE)
        txt(s, desc,  Inches(1.6), y + Inches(0.45), Inches(9), Inches(0.6),
            size=Pt(15), color=RGBColor(0xAA, 0xAA, 0xCC))


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 – Roadmap
# ─────────────────────────────────────────────────────────────────────────────
def slide_roadmap(p):
    s = blank(p)
    bg(s, LIGHT_BG)
    box(s, Inches(0), Inches(0), SLIDE_W, Inches(1.3), fill_color=PURPLE)
    txt(s, "Nächste Schritte", Inches(0.5), Inches(0.25), Inches(10), Inches(0.8),
        size=Pt(32), bold=True, color=WHITE)

    done = [
        "✅ IB-CSV Parser (automatische Erkennung)",
        "✅ FX-Umrechnung USD → EUR",
        "✅ Asset-Typ-Erkennung (Aktie, ETF, Option, Anleihe …)",
        "✅ Konzentrations-Warnungen",
        "✅ KI-Kommentar (deutsch)",
    ]
    next_ = [
        "🔲 PDF-Parser verbessern (Tabellen-Erkennung)",
        "🔲 Historische Performance (Rendite-Chart)",
        "🔲 Mehrere Depots vergleichen",
        "🔲 Export als PDF-Report",
        "🔲 AWS Deployment abschließen",
    ]

    txt(s, "Fertig", Inches(0.5), Inches(1.5), Inches(5.5), Inches(0.5),
        size=Pt(18), bold=True, color=GREEN)
    bullet_box(s, done, Inches(0.5), Inches(2.0), Inches(5.8), Inches(4.5),
               size=Pt(15), color=DARK, bullet="")

    txt(s, "Geplant", Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.5),
        size=Pt(18), bold=True, color=PURPLE)
    bullet_box(s, next_, Inches(7.0), Inches(2.0), Inches(5.8), Inches(4.5),
               size=Pt(15), color=DARK, bullet="")

    # Divider
    box(s, Inches(6.6), Inches(1.5), Inches(0.04), Inches(5.5),
        fill_color=RGBColor(0xCC, 0xCC, 0xDD))


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 – End / Contact
# ─────────────────────────────────────────────────────────────────────────────
def slide_end(p):
    s = blank(p)
    bg(s, DARK)
    box(s, Inches(0), Inches(0), Inches(0.18), SLIDE_H, fill_color=BLUE)

    txt(s, "Depot-MVP",
        Inches(0.5), Inches(2.0), Inches(12), Inches(1.2),
        size=Pt(54), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    txt(s, "Portfolio-Analyse · sofort · automatisch · kostenlos",
        Inches(0.5), Inches(3.3), Inches(12), Inches(0.7),
        size=Pt(22), color=BLUE, align=PP_ALIGN.CENTER)

    txt(s, "Entwickelt von Tim Solovyev  ·  Dr. Mayer Consulting GmbH",
        Inches(0.5), Inches(5.2), Inches(12), Inches(0.5),
        size=Pt(18), bold=True, color=RGBColor(0xBB, 0xBB, 0xDD), align=PP_ALIGN.CENTER)

    txt(s, "github  ·  AWS EB  ·  localhost:8080",
        Inches(0.5), Inches(5.8), Inches(12), Inches(0.5),
        size=Pt(16), color=RGBColor(0x88, 0x88, 0xAA), align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# Build & save
# ─────────────────────────────────────────────────────────────────────────────
def main():
    p = prs()
    slide_title(p)
    slide_problem(p)
    slide_solution(p)
    slide_output(p)
    slide_analysis(p)
    slide_tech(p)
    slide_demo(p)
    slide_roadmap(p)
    slide_end(p)

    out = "Depot-MVP Präsentation.pptx"
    p.save(out)
    print(f"Gespeichert: {out}")


if __name__ == "__main__":
    main()
